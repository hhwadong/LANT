#!/usr/bin/env python3
import os
import json
import glob
import shutil
import ollama
from datetime import datetime
import PyPDF2
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image
import markdown
import re
import hashlib
import time
import psutil

# Configuration constants for scalability
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB max file size
MEMORY_WARNING_THRESHOLD = 80  # Warn when memory usage > 80%
MAX_CONTEXT_LENGTH = 12000  # Max characters for AI context
CHUNK_SIZE = 10  # Pages/slides to process at once

class SilentDirectoryAssistant:
    def __init__(self, model="codellama:7b"):
        self.model = model
        self.base_dir = "learning_assistant"
        self.lectures_dir = os.path.join(self.base_dir, "lectures")
        self.cache_dir = os.path.join(self.base_dir, "cache")
        self.current_lecture = None
        self.current_session = None
        # Default model parameters
        self.model_params = {
            "temperature": 0.7,
            "top_p": 0.9,
            "num_predict": 3072
        }
        self.max_context_messages = 12  # Maximum messages before summarization
        self.ensure_directories_silent()
    
    def ensure_directories_silent(self):
        """Create necessary directories silently if they don't exist"""
        directories = [
            self.base_dir,
            self.lectures_dir,
            self.cache_dir
        ]
        
        for directory in directories:
            if not os.path.exists(directory):
                os.makedirs(directory)
    
    def get_file_hash(self, file_path):
        """Generate a hash for a file based on its path and modification time"""
        stat = os.stat(file_path)
        hash_input = f"{file_path}_{stat.st_mtime}_{stat.st_size}"
        return hashlib.md5(hash_input.encode()).hexdigest()

    def check_memory_usage(self):
        """Check current memory usage and warn if high"""
        try:
            memory = psutil.virtual_memory()
            if memory.percent > MEMORY_WARNING_THRESHOLD:
                print(f"⚠️  Warning: High memory usage ({memory.percent:.1f}%)")
            return memory.percent
        except Exception:
            # If psutil fails, continue silently
            return None

    def check_file_size(self, file_path):
        """Check if file is too large and warn if necessary"""
        try:
            size = os.path.getsize(file_path)
            if size > MAX_FILE_SIZE:
                size_mb = size / (1024 * 1024)
                print(f"⚠️  Warning: Large file ({size_mb:.1f}MB) - processing may be slow")
            return size
        except OSError:
            return None
    
    def get_cached_text(self, file_path):
        """Get cached text for a document if available and valid"""
        if not os.path.exists(file_path):
            return None
            
        file_hash = self.get_file_hash(file_path)
        cache_file = os.path.join(self.cache_dir, f"{file_hash}.txt")
        
        if os.path.exists(cache_file):
            try:
                with open(cache_file, 'r', encoding='utf-8') as f:
                    return f.read()
            except (FileNotFoundError, UnicodeDecodeError, PermissionError, OSError) as e:
                # Silently return None for cache errors - cache is optional
                return None
        return None
    
    def cache_text(self, file_path, text):
        """Cache extracted text for a document"""
        file_hash = self.get_file_hash(file_path)
        cache_file = os.path.join(self.cache_dir, f"{file_hash}.txt")
        
        try:
            with open(cache_file, 'w', encoding='utf-8') as f:
                f.write(text)
            return True
        except (PermissionError, OSError, UnicodeEncodeError) as e:
            # Silently return False for cache errors - cache is optional
            return False
    
    def list_lectures(self):
        """List all lectures"""
        lectures = []
        if os.path.exists(self.lectures_dir):
            for item in os.listdir(self.lectures_dir):
                if os.path.isdir(os.path.join(self.lectures_dir, item)):
                    lectures.append(item)
        return sorted(lectures)
    
    def create_lecture(self, lecture_name):
        """Create a new lecture with proper directory structure"""
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        
        if os.path.exists(lecture_path):
            return f"Lecture '{lecture_name}' already exists"
        
        # Create lecture directory structure
        os.makedirs(lecture_path)
        os.makedirs(os.path.join(lecture_path, "sessions"))
        os.makedirs(os.path.join(lecture_path, "docs"))
        
        # Create lecture_info.json
        lecture_info = {
            "name": lecture_name,
            "created_at": datetime.now().isoformat(),
            "documents": [],  # List of documents available to all sessions
            "sessions": [],
            "current_model": self.model,
            "model_params": self.model_params.copy(),
            "directory_structure": {
                "root": lecture_path,
                "sessions": os.path.join(lecture_path, "sessions"),
                "docs": os.path.join(lecture_path, "docs")
            }
        }
        
        with open(os.path.join(lecture_path, "lecture_info.json"), 'w') as f:
            json.dump(lecture_info, f, indent=2)
        
        return f"Created lecture: {lecture_name}"
    
    def add_lecture(self, lecture_name):
        """Add a new lecture by name"""
        return self.create_lecture(lecture_name)
    
    def load_lecture(self, lecture_name):
        """Load a lecture"""
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        
        if not os.path.exists(lecture_path):
            return False
        
        self.current_lecture = lecture_name
        
        # Load lecture info
        with open(os.path.join(lecture_path, "lecture_info.json"), 'r') as f:
            lecture_info = json.load(f)
            self.model = lecture_info.get("current_model", self.model)
            self.model_params = lecture_info.get("model_params", self.model_params.copy())
        
        return True
    
    def get_lecture_info(self, lecture_name=None):
        """Get lecture information"""
        if not lecture_name:
            lecture_name = self.current_lecture
        
        if not lecture_name:
            return None
        
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        if not os.path.exists(lecture_path):
            return None
        
        with open(os.path.join(lecture_path, "lecture_info.json"), 'r') as f:
            return json.load(f)
    
    def list_sessions(self, lecture_name=None):
        """List all sessions in a lecture"""
        if not lecture_name:
            lecture_name = self.current_lecture
        
        if not lecture_name:
            return []
        
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        sessions_path = os.path.join(lecture_path, "sessions")
        
        if not os.path.exists(sessions_path):
            return []
        
        sessions = []
        for file in os.listdir(sessions_path):
            if file.endswith('.json'):
                sessions.append(os.path.splitext(file)[0])
        
        return sorted(sessions)
    
    def create_session(self, session_name=None, lecture_name=None):
        """Create a new session in a lecture"""
        if not lecture_name:
            lecture_name = self.current_lecture
        
        if not lecture_name:
            return "No lecture selected"
        
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        sessions_path = os.path.join(lecture_path, "sessions")
        
        if not session_name:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            session_name = f"session_{timestamp}"
        
        session_path = os.path.join(sessions_path, f"{session_name}.json")
        
        if os.path.exists(session_path):
            return f"Session '{session_name}' already exists in lecture '{lecture_name}'"
        
        # Create session file
        session_data = {
            "name": session_name,
            "lecture": lecture_name,
            "created_at": datetime.now().isoformat(),
            "model": self.model,
            "model_params": self.model_params.copy(),
            "messages": [],
            "documents": [],  # List of documents specific to this session
            "lectures_referenced": [],
            "summaries": [],  # List of conversation summaries
            "file_path": session_path
        }
        
        with open(session_path, 'w') as f:
            json.dump(session_data, f, indent=2)
        
        # Update lecture info
        self.update_lecture_sessions(lecture_name)
        
        self.current_session = session_name
        return f"Created session '{session_name}' in lecture '{lecture_name}'"
    
    def load_session(self, session_name, lecture_name=None):
        """Load a session"""
        if not lecture_name:
            lecture_name = self.current_lecture
        
        if not lecture_name:
            return False
        
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{session_name}.json")
        
        if not os.path.exists(session_path):
            return False
        
        self.current_session = session_name
        
        # Load session data (lazy loading - only metadata)
        with open(session_path, 'r') as f:
            session_data = json.load(f)
            self.model = session_data.get("model", self.model)
            self.model_params = session_data.get("model_params", self.model_params.copy())
        
        return True
    
    def get_session_history(self, limit=None, offset=0):
        """Get conversation history of current session with lazy loading"""
        if not self.current_lecture or not self.current_session:
            return []
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")
        
        try:
            with open(session_path, 'r', encoding='utf-8') as f:
                session_data = json.load(f)
                messages = session_data.get("messages", [])

                # Apply lazy loading parameters
                if limit is not None:
                    end = offset + limit
                    return messages[offset:end]
                return messages
        except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError) as e:
            # Return empty list for session reading errors
            return []
    
    def save_message(self, role, content, lecture_ref=None):
        """Save a message to the current session"""
        if not self.current_lecture or not self.current_session:
            return
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")
        
        try:
            with open(session_path, 'r', encoding='utf-8') as f:
                session_data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError):
            # Create new session data if file doesn't exist or is corrupted
            session_data = {
                "name": self.current_session,
                "messages": [],
                "lectures_referenced": [],
                "model": self.model,
                "model_params": self.model_params.copy(),
                "summaries": []
            }
        
        message = {
            "role": role,
            "content": content,
            "timestamp": datetime.now().isoformat()
        }
        
        # Add model information for assistant messages
        if role == "assistant":
            message["model"] = self.model
            message["model_params"] = self.model_params.copy()
        
        if lecture_ref:
            message["lecture_ref"] = lecture_ref
            if lecture_ref not in session_data.get("lectures_referenced", []):
                session_data["lectures_referenced"].append(lecture_ref)
        
        session_data["messages"].append(message)
        
        with open(session_path, 'w') as f:
            json.dump(session_data, f, indent=2)

    def get_conversation_history(self):
        """Get all messages for the current session"""
        return self.get_session_history()

    def clear_conversation_history(self):
        """Clear all messages for the current session"""
        if not self.current_lecture or not self.current_session:
            return False

        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")

        try:
            with open(session_path, 'r', encoding='utf-8') as f:
                session_data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError):
            # Create new session data if file doesn't exist or is corrupted
            session_data = {
                "name": self.current_session,
                "messages": [],
                "lectures_referenced": [],
                "model": self.model,
                "model_params": self.model_params.copy(),
                "summaries": []
            }

        # Clear messages but keep other data
        session_data["messages"] = []
        session_data["summaries"] = []  # Also clear summaries

        with open(session_path, 'w') as f:
            json.dump(session_data, f, indent=2)

        return True

    def update_lecture_sessions(self, lecture_name):
        """Update lecture's session list"""
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        if not os.path.exists(lecture_path):
            return
        
        sessions = self.list_sessions(lecture_name)
        
        with open(os.path.join(lecture_path, "lecture_info.json"), 'r+') as f:
            lecture_info = json.load(f)
            lecture_info["sessions"] = sessions
            f.seek(0)
            json.dump(lecture_info, f, indent=2)
            f.truncate()
    
    def get_status(self):
        """Get status of all lectures and sessions"""
        status = {
            "directory_structure": {
                "base": self.base_dir,
                "lectures": self.lectures_dir,
                "cache": self.cache_dir
            },
            "lectures": {},
            "total_sessions": 0,
            "total_estimated_tokens": 0,
            "cache_stats": self.get_cache_stats()
        }
        
        if os.path.exists(self.lectures_dir):
            for lecture_name in self.list_lectures():
                lecture_path = os.path.join(self.lectures_dir, lecture_name)
                lecture_info = self.get_lecture_info(lecture_name)
                
                if lecture_info:
                    lecture_status = {
                        "name": lecture_name,
                        "created_at": lecture_info.get("created_at", "Unknown"),
                        "documents": lecture_info.get("documents", []),
                        "directory": lecture_path,
                        "sessions": {},
                        "session_count": 0,
                        "estimated_tokens": 0
                    }
                    
                    # Process each session
                    for session_name in self.list_sessions(lecture_name):
                        session_path = os.path.join(lecture_path, "sessions", f"{session_name}.json")
                        try:
                            with open(session_path, 'r') as f:
                                session_data = json.load(f)
                                
                                # Calculate estimated tokens
                                messages = session_data.get("messages", [])
                                total_chars = sum(len(msg.get('content', '')) for msg in messages)
                                estimated_tokens = total_chars // 4
                                
                                session_status = {
                                    "name": session_name,
                                    "created_at": session_data.get("created_at", "Unknown"),
                                    "model": session_data.get("model", "Unknown"),
                                    "message_count": len(messages),
                                    "estimated_tokens": estimated_tokens,
                                    "documents": session_data.get("documents", []),
                                    "lectures_referenced": session_data.get("lectures_referenced", []),
                                    "file_path": session_path,
                                    "summaries": session_data.get("summaries", [])
                                }
                                
                                lecture_status["sessions"][session_name] = session_status
                                lecture_status["session_count"] += 1
                                lecture_status["estimated_tokens"] += estimated_tokens
                                
                        except Exception as e:
                            lecture_status["sessions"][session_name] = {"error": str(e)}
                    
                    status["lectures"][lecture_name] = lecture_status
                    status["total_sessions"] += lecture_status["session_count"]
                    status["total_estimated_tokens"] += lecture_status["estimated_tokens"]
        
        return status
    
    def get_cache_stats(self):
        """Get statistics about the cache"""
        stats = {
            "cache_dir": self.cache_dir,
            "cache_files": 0,
            "cache_size_mb": 0
        }
        
        if os.path.exists(self.cache_dir):
            cache_files = os.listdir(self.cache_dir)
            stats["cache_files"] = len(cache_files)
            
            total_size = 0
            for file in cache_files:
                file_path = os.path.join(self.cache_dir, file)
                if os.path.isfile(file_path):
                    total_size += os.path.getsize(file_path)
            
            stats["cache_size_mb"] = round(total_size / (1024 * 1024), 2)
        
        return stats
    
    def clear_cache(self):
        """Clear the document cache"""
        if os.path.exists(self.cache_dir):
            for file in os.listdir(self.cache_dir):
                file_path = os.path.join(self.cache_dir, file)
                if os.path.isfile(file_path):
                    os.remove(file_path)
        return "Cache cleared"
    
    def set_model_parameter(self, param, value):
        """Set a model parameter for the current session/lecture"""
        try:
            if param == "temperature":
                value = float(value)
                if not 0.0 <= value <= 1.0:
                    return "Temperature must be between 0.0 and 1.0"
            elif param == "top_p":
                value = float(value)
                if not 0.0 <= value <= 1.0:
                    return "Top-p must be between 0.0 and 1.0"
            elif param == "num_predict":
                value = int(value)
                if value < 1:
                    return "Num predict must be a positive integer"
            else:
                return f"Unknown parameter: {param}"
            
            self.model_params[param] = value
            
            # Update session data if a session is active
            if self.current_lecture and self.current_session:
                lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
                sessions_path = os.path.join(lecture_path, "sessions")
                session_path = os.path.join(sessions_path, f"{self.current_session}.json")
                
                try:
                    with open(session_path, 'r') as f:
                        session_data = json.load(f)
                    
                    session_data["model_params"] = self.model_params.copy()
                    
                    with open(session_path, 'w') as f:
                        json.dump(session_data, f, indent=2)
                except (PermissionError, OSError):
                    pass
            
            return f"Set {param} to {value}"
        except ValueError:
            return f"Invalid value for {param}: {value}"
    
    def list_model_parameters(self):
        """List current model parameters"""
        params = self.model_params.copy()
        params["model"] = self.model
        return params
    
    def summarize_conversation(self, messages):
        """Generate a summary of the conversation"""
        if not messages:
            return ""
        
        # Create a summary prompt
        summary_prompt = f"""
Please provide a concise summary of the following conversation between a user and an AI assistant. 
Focus on the key topics discussed, important questions asked, and main points from the responses.
The summary should be brief but comprehensive enough to provide context for continuing the conversation.

CONVERSATION:
---
{json.dumps(messages, indent=2)}
---

SUMMARY:
"""
        
        try:
            # Use a more conservative temperature for summarization
            summary_params = self.model_params.copy()
            summary_params["temperature"] = 0.3
            summary_params["num_predict"] = 1024
            
            response = ollama.chat(
                model=self.model,
                messages=[{'role': 'user', 'content': summary_prompt}],
                options=summary_params
            )
            
            return response['message']['content'].strip()
        except Exception as e:
            return f"Error generating summary: {str(e)}"
    
    def get_conversation_context(self):
        """Get conversation context with summarization for long sessions"""
        if not self.current_lecture or not self.current_session:
            return []
        
        # Get all messages
        messages = self.get_session_history()
        
        # If we have few messages, return them directly
        if len(messages) <= self.max_context_messages:
            return messages
        
        # For long conversations, use the most recent messages and a summary of earlier ones
        recent_messages = messages[-self.max_context_messages:]
        earlier_messages = messages[:-self.max_context_messages]
        
        # Generate summary of earlier messages
        summary = self.summarize_conversation(earlier_messages)
        
        # Create a context message with the summary
        context_message = {
            "role": "system",
            "content": f"Earlier in this conversation: {summary}",
            "timestamp": datetime.now().isoformat()
        }
        
        # Save the summary to the session
        self.save_summary(summary, 0, len(earlier_messages) - 1)
        
        # Return the context message followed by recent messages
        return [context_message] + recent_messages
    
    def save_summary(self, summary, start_index, end_index):
        """Save a conversation summary to the current session"""
        if not self.current_lecture or not self.current_session:
            return
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")
        
        try:
            with open(session_path, 'r', encoding='utf-8') as f:
                session_data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError):
            session_data = {"name": self.current_session, "summaries": []}
        
        summary_data = {
            "summary": summary,
            "start_index": start_index,
            "end_index": end_index,
            "timestamp": datetime.now().isoformat()
        }
        
        if "summaries" not in session_data:
            session_data["summaries"] = []
        
        session_data["summaries"].append(summary_data)
        
        with open(session_path, 'w') as f:
            json.dump(session_data, f, indent=2)
    
    def list_summaries(self):
        """List conversation summaries for the current session"""
        if not self.current_lecture or not self.current_session:
            return []
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")
        
        try:
            with open(session_path, 'r') as f:
                session_data = json.load(f)
                return session_data.get("summaries", [])
        except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError):
            return []
    
    def extract_pdf_text(self, pdf_path):
        """Extract text from PDF file with enhanced formatting"""
        # Check file size first
        self.check_file_size(pdf_path)
        self.check_memory_usage()

        # Check cache first
        cached_text = self.get_cached_text(pdf_path)
        if cached_text is not None:
            return cached_text

        text = ""
        try:
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                # Extract text with page numbers for better context
                text_parts = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():  # Only add if there's actual text
                        text_parts.append(f"\n--- Page {i+1} ---\n")
                        text_parts.append(page_text)
                        text_parts.append("\n")

                # Join all parts efficiently
                text += "".join(text_parts)
                
                # Try to extract images from PDF and run OCR
                try:
                    for i, page in enumerate(reader.pages):
                        if '/Resources' in page and '/XObject' in page['/Resources']:
                            xObject = page['/Resources']['/XObject'].get_object()
                            for obj in xObject:
                                if xObject[obj]['/Subtype'] == '/Image':
                                    try:
                                        # Extract image data
                                        data = xObject[obj]._data
                                        # Create a temporary image file
                                        temp_img = f"temp_img_{i}_{obj[1:]}.png"
                                        with open(temp_img, "wb") as img_file:
                                            img_file.write(data)
                                        
                                        # Run OCR on the image
                                        img_text = self.extract_image_text(temp_img)
                                        if img_text and not img_text.startswith("Error"):
                                            text += f"\n--- OCR from Page {i+1} Image ---\n"
                                            text += img_text + "\n"
                                        
                                        # Clean up temp file
                                        os.remove(temp_img)
                                    except (OSError, PermissionError, Exception):
                                        # Continue if OCR image processing fails
                                        continue
                except Exception:
                    # Continue if PDF image extraction fails
                    pass
                
            # Cache the extracted text
            self.cache_text(pdf_path, text)
            return text
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"
    
    def extract_ppt_text(self, ppt_path):
        """Extract text from PowerPoint file with enhanced formatting"""
        # Check cache first
        cached_text = self.get_cached_text(ppt_path)
        if cached_text is not None:
            return cached_text
            
        text = ""
        try:
            prs = Presentation(ppt_path)
            
            # Extract text with slide numbers and titles
            text_parts = []
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"\n--- Slide {i+1} ---\n")

                # Try to get slide title
                if slide.shapes.title:
                    text_parts.append(f"Title: {slide.shapes.title.text}\n")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text.strip()
                        if shape_text:
                            text_parts.append(shape_text)
                            text_parts.append("\n")

            # Join all parts efficiently
            text = "".join(text_parts)

            # Try to extract images from slides and run OCR
            try:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Shape type for pictures
                        try:
                            # Get image data
                            image = shape.image
                            # Create a temporary image file
                            temp_img = f"temp_img_{i}_{shape.shape_id}.png"
                            with open(temp_img, "wb") as img_file:
                                img_file.write(image.blob)

                            # Run OCR on the image
                            img_text = self.extract_image_text(temp_img)
                            if img_text and not img_text.startswith("Error"):
                                text += f"\n--- OCR from Slide {i+1} Image ---\n"
                                text += img_text + "\n"

                            # Clean up temp file
                            os.remove(temp_img)
                        except (OSError, PermissionError, Exception):
                            continue
            except Exception:
                pass
                
            # Cache the extracted text
            self.cache_text(ppt_path, text)
            return text
        except Exception as e:
            return f"Error extracting PowerPoint text: {str(e)}"
    
    def extract_docx_text(self, docx_path):
        """Extract text from Word document with enhanced formatting"""
        # Check cache first
        cached_text = self.get_cached_text(docx_path)
        if cached_text is not None:
            return cached_text
            
        text = ""
        try:
            doc = Document(docx_path)
            
            # Extract document properties
            core_props = doc.core_properties
            if core_props.title:
                text += f"Title: {core_props.title}\n"
            if core_props.author:
                text += f"Author: {core_props.author}\n"
            if core_props.subject:
                text += f"Subject: {core_props.subject}\n"
            text += "\n"
            
            # Extract paragraphs with style information
            for para in doc.paragraphs:
                if para.text.strip():
                    style_name = para.style.name.lower()
                    
                    # Add formatting based on style
                    if "heading" in style_name:
                        # Extract heading level
                        level = re.search(r'heading (\d+)', style_name)
                        if level:
                            heading_level = int(level.group(1))
                            text += f"\n{'#' * heading_level} {para.text}\n\n"
                        else:
                            text += f"\n{para.text}\n\n"
                    elif "title" in style_name:
                        text += f"\n# {para.text}\n\n"
                    else:
                        text += para.text + "\n\n"
            
            # Extract tables
            for i, table in enumerate(doc.tables):
                text += f"\n--- Table {i+1} ---\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text += " | ".join(row_text) + "\n"
                text += "\n"
            
            # Try to extract images and run OCR
            try:
                rels = doc.part.rels
                for rel in rels:
                    if "image" in rels[rel].target_ref:
                        try:
                            # Get image data
                            image_part = rels[rel].target_part
                            # Create a temporary image file
                            temp_img = f"temp_img_{rel}.png"
                            with open(temp_img, "wb") as img_file:
                                img_file.write(image_part.blob)
                            
                            # Run OCR on the image
                            img_text = self.extract_image_text(temp_img)
                            if img_text and not img_text.startswith("Error"):
                                text += f"\n--- OCR from Document Image ---\n"
                                text += img_text + "\n"
                            
                            # Clean up temp file
                            os.remove(temp_img)
                        except (OSError, PermissionError, Exception):
                            continue
            except Exception:
                pass
            
            # Cache the extracted text
            self.cache_text(docx_path, text)
            return text
        except Exception as e:
            return f"Error extracting Word document text: {str(e)}"
    
    def extract_txt_text(self, txt_path):
        """Extract text from a plain text file"""
        # Check cache first
        cached_text = self.get_cached_text(txt_path)
        if cached_text is not None:
            return cached_text
            
        try:
            with open(txt_path, 'r', encoding='utf-8') as file:
                text = file.read()
                # Cache the extracted text
                self.cache_text(txt_path, text)
                return text
        except Exception as e:
            return f"Error reading text file: {str(e)}"
    
    def extract_md_text(self, md_path):
        """Extract text from a Markdown file with enhanced formatting"""
        # Check cache first
        cached_text = self.get_cached_text(md_path)
        if cached_text is not None:
            return cached_text
            
        try:
            with open(md_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
                
                # Convert markdown to plain text while preserving structure
                html = markdown.markdown(md_content)
                
                # Simple HTML to text conversion
                text = re.sub(r'<[^>]+>', '', html)  # Remove HTML tags
                
                # Add markdown structure indicators
                lines = md_content.split('\n')
                structured_text = ""
                
                for line in lines:
                    # Headers
                    if line.startswith('# '):
                        structured_text += f"\n# {line[2:]}\n\n"
                    elif line.startswith('## '):
                        structured_text += f"\n## {line[3:]}\n\n"
                    elif line.startswith('### '):
                        structured_text += f"\n### {line[4:]}\n\n"
                    # Lists
                    elif line.startswith('- ') or line.startswith('* '):
                        structured_text += f"- {line[2:]}\n"
                    # Code blocks
                    elif line.startswith('```'):
                        structured_text += "\n--- Code Block ---\n"
                    # Regular text
                    elif line.strip():
                        structured_text += line + "\n"
                    # Empty lines
                    else:
                        structured_text += "\n"
                
                # Cache the extracted text
                self.cache_text(md_path, structured_text)
                return structured_text
        except Exception as e:
            return f"Error reading Markdown file: {str(e)}"
    
    def extract_image_text(self, image_path):
        """Extract text from an image using OCR"""
        # Check cache first
        cached_text = self.get_cached_text(image_path)
        if cached_text is not None:
            return cached_text
            
        try:
            # Open the image
            img = Image.open(image_path)
            
            # Preprocess the image for better OCR
            # Convert to grayscale
            img = img.convert('L')
            
            # Use pytesseract to extract text
            text = pytesseract.image_to_string(img)
            
            # Cache the extracted text
            self.cache_text(image_path, text)
            return text
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    def extract_document_text(self, file_path):
        """Extract text from various document types with enhanced formatting"""
        if not os.path.exists(file_path):
            return f"File not found: {file_path}"
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self.extract_pdf_text(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            return self.extract_ppt_text(file_path)
        elif file_ext == '.docx':
            return self.extract_docx_text(file_path)
        elif file_ext == '.txt':
            return self.extract_txt_text(file_path)
        elif file_ext == '.md':
            return self.extract_md_text(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']:
            return self.extract_image_text(file_path)
        else:
            return f"Unsupported document type: {file_ext}"
    
    def add_document_to_lecture(self, file_path):
        """Add a document to the current lecture (available to all sessions)"""
        if not self.current_lecture:
            return "No lecture selected"
        
        if not os.path.exists(file_path):
            return f"File not found: {file_path}"
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        docs_path = os.path.join(lecture_path, "docs")
        
        # Copy file to docs directory
        dest_path = os.path.join(docs_path, os.path.basename(file_path))
        shutil.copy2(file_path, dest_path)
        
        # Update lecture info
        lecture_info_path = os.path.join(lecture_path, "lecture_info.json")
        with open(lecture_info_path, 'r') as f:
            lecture_info = json.load(f)
        
        # Add document to the list if not already there
        doc_name = os.path.basename(file_path)
        if doc_name not in lecture_info.get("documents", []):
            lecture_info["documents"].append(doc_name)
        
        with open(lecture_info_path, 'w') as f:
            json.dump(lecture_info, f, indent=2)
        
        return f"Added document '{doc_name}' to lecture '{self.current_lecture}' (available to all sessions)"
    
    def add_document_to_session(self, file_path):
        """Add a document to the current session (available only to this session)"""
        if not self.current_lecture or not self.current_session:
            return "No lecture or session selected"
        
        if not os.path.exists(file_path):
            return f"File not found: {file_path}"
        
        lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
        sessions_path = os.path.join(lecture_path, "sessions")
        session_path = os.path.join(sessions_path, f"{self.current_session}.json")
        
        # Create session-specific docs directory if it doesn't exist
        session_docs_dir = os.path.join(lecture_path, "session_docs", self.current_session)
        if not os.path.exists(session_docs_dir):
            os.makedirs(session_docs_dir)
        
        # Copy file to session-specific docs directory
        dest_path = os.path.join(session_docs_dir, os.path.basename(file_path))
        shutil.copy2(file_path, dest_path)
        
        # Update session info
        with open(session_path, 'r') as f:
            session_data = json.load(f)
        
        # Add document to the list if not already there
        doc_name = os.path.basename(file_path)
        if doc_name not in session_data.get("documents", []):
            session_data["documents"].append(doc_name)
        
        with open(session_path, 'w') as f:
            json.dump(session_data, f, indent=2)
        
        return f"Added document '{doc_name}' to session '{self.current_session}' (available only to this session)"
    
    def analyze_lecture(self, lecture_name, question):
        """Analyze a specific lecture document"""
        if not self.current_lecture:
            return "No lecture selected"
        
        lecture_path = os.path.join(self.lectures_dir, lecture_name)
        if not os.path.exists(lecture_path):
            return f"Lecture not found: {lecture_name}"
        
        # Get lecture info
        lecture_info = self.get_lecture_info(lecture_name)
        if not lecture_info:
            return f"Could not load lecture info for: {lecture_name}"
        
        # Collect documents to analyze
        docs_to_analyze = []
        
        # If we're in a session, include session-specific documents
        if self.current_session:
            session_path = os.path.join(lecture_path, "sessions", f"{self.current_session}.json")
            if os.path.exists(session_path):
                with open(session_path, 'r') as f:
                    session_data = json.load(f)
                    session_docs = session_data.get("documents", [])
                    
                    # Add session-specific documents
                    for doc in session_docs:
                        session_docs_dir = os.path.join(lecture_path, "session_docs", self.current_session)
                        doc_path = os.path.join(session_docs_dir, doc)
                        if os.path.exists(doc_path):
                            docs_to_analyze.append(doc_path)
        
        # Add lecture-level documents (available to all sessions)
        lecture_docs = lecture_info.get("documents", [])
        docs_path = os.path.join(lecture_path, "docs")
        
        for doc in lecture_docs:
            doc_path = os.path.join(docs_path, doc)
            if os.path.exists(doc_path):
                docs_to_analyze.append(doc_path)
        
        if not docs_to_analyze:
            return "No documents found for analysis"

        # Check memory usage before processing multiple documents
        self.check_memory_usage()

        # Check file sizes for all documents
        for doc_path in docs_to_analyze:
            self.check_file_size(doc_path)

        # Extract text from all documents
        text_parts = []
        for doc_path in docs_to_analyze:
            doc_name = os.path.basename(doc_path)
            extracted_text = self.extract_document_text(doc_path)

            # Only add if extraction was successful
            if not extracted_text.startswith("Error"):
                text_parts.append(f"\n--- Document: {doc_name} ---\n")
                text_parts.append(extracted_text)
                text_parts.append("\n")
            else:
                text_parts.append(f"\n--- Error processing document: {doc_name} ---\n")
                text_parts.append(extracted_text)
                text_parts.append("\n")

        # Join all parts efficiently
        all_text = "".join(text_parts)
        
        # Create prompt
        prompt = f"""
I'm analyzing lecture documents for my computer science studies. 

LECTURE CONTENT:
---
{all_text[:MAX_CONTEXT_LENGTH]}  # Use configurable limit
---

MY QUESTION:
{question}

Please provide a comprehensive analysis that includes:
1. Direct answer to my question
2. Explanation of relevant concepts
3. Connection to broader computer science topics
4. Any examples or applications mentioned

Be thorough and detailed in your response.
"""
        
        # Get response
        try:
            response = ollama.chat(
                model=self.model,
                messages=[{'role': 'user', 'content': prompt}],
                options=self.model_params
            )
            
            # Save to session
            self.save_message("user", f"Analyzed lecture: {lecture_name}\nQuestion: {question}", lecture_name)
            self.save_message("assistant", response['message']['content'], lecture_name)
            
            return response['message']['content']
            
        except Exception as e:
            return f"Error getting response: {str(e)}"
    
    def merge_all_sessions(self, lecture_name=None):
        """Merge all sessions in a lecture into a new session with batch processing"""
        if not lecture_name:
            lecture_name = self.current_lecture
        
        if not lecture_name:
            return "No lecture selected"
        
        # Get all sessions in the lecture
        sessions = self.list_sessions(lecture_name)
        if not sessions:
            return f"No sessions found in lecture '{lecture_name}'"
        
        print(f"Starting merge of {len(sessions)} sessions...")
        print("This may take a moment for large sessions...")
        print()
        
        # Check estimated size
        total_estimated_size = 0
        for session_name in sessions:
            session_path = os.path.join(self.lectures_dir, lecture_name, "sessions", f"{session_name}.json")
            try:
                with open(session_path, 'r') as f:
                    session_data = json.load(f)
                    messages = session_data.get("messages", [])
                    # Rough estimate: 1KB per message
                    total_estimated_size += len(messages) * 1024
            except (FileNotFoundError, json.JSONDecodeError, UnicodeDecodeError, PermissionError):
                continue
        
        # Warn if large
        if total_estimated_size > 10 * 1024 * 1024:  # 10MB
            print(f"Warning: This merge may create a large session (~{total_estimated_size / (1024*1024):.1f} MB)")
            confirm = input("Continue? (y/n): ")
            if confirm.lower() != 'y':
                return "Merge cancelled by user"
        
        # Create merged session
        merged_session_name = f"MergedSess-{lecture_name}"
        sessions_path = os.path.join(self.lectures_dir, lecture_name, "sessions")
        merged_session_path = os.path.join(sessions_path, f"{merged_session_name}.json")
        
        # Check if already exists
        if os.path.exists(merged_session_path):
            # Remove existing merged session if user confirms
            confirm = input(f"Merged session '{merged_session_name}' already exists. Overwrite? (y/n): ")
            if confirm.lower() != 'y':
                return "Merge cancelled"
            os.remove(merged_session_path)
        
        # Initialize merged session data
        merged_session_data = {
            "name": merged_session_name,
            "lecture": lecture_name,
            "created_at": datetime.now().isoformat(),
            "model": self.model,
            "model_params": self.model_params.copy(),
            "messages": [],
            "documents": [],  # Merged session doesn't inherit documents
            "lectures_referenced": [],
            "summaries": [],
            "merge_info": {
                "original_sessions": sessions,
                "total_messages": 0,
                "merged_at": datetime.now().isoformat(),
                "estimated_size_mb": total_estimated_size / (1024 * 1024),
                "file_path": merged_session_path
            }
        }
        
        # Process sessions in batches
        batch_size = 5  # Process 5 sessions at a time
        total_messages = 0
        lecture_refs = set()
        
        for i in range(0, len(sessions), batch_size):
            batch_sessions = sessions[i:i+batch_size]
            print(f"Processing batch {i//batch_size + 1}/{(len(sessions) + batch_size - 1)//batch_size}...")
            
            for session_name in batch_sessions:
                session_path = os.path.join(self.lectures_dir, lecture_name, "sessions", f"{session_name}.json")
                try:
                    with open(session_path, 'r') as f:
                        session_data = json.load(f)
                        messages = session_data.get("messages", [])
                        
                        if messages:
                            # Add session header
                            merged_session_data["messages"].append({
                                "role": "system",
                                "content": f"=== Session: {session_name} ===",
                                "timestamp": datetime.now().isoformat()
                            })
                            
                            # Add all messages from this session
                            merged_session_data["messages"].extend(messages)
                            total_messages += len(messages)
                            
                            # Add session footer
                            merged_session_data["messages"].append({
                                "role": "system",
                                "content": f"=== End of Session: {session_name} ===",
                                "timestamp": datetime.now().isoformat()
                            })
                            
                            # Collect lecture references
                            for msg in messages:
                                if msg.get("role") in ["user", "assistant"]:
                                    ref = msg.get("lecture_ref")
                                    if ref:
                                        lecture_refs.add(ref)
                            
                            print(f"  - Processed session '{session_name}' ({len(messages)} messages)")
                            
                except Exception as e:
                    print(f"  - Error processing session '{session_name}': {e}")
                    continue
            
            # Save progress after each batch
            merged_session_data["merge_info"]["total_messages"] = total_messages
            merged_session_data["lectures_referenced"] = list(lecture_refs)
            
            with open(merged_session_path, 'w') as f:
                json.dump(merged_session_data, f, indent=2)
            
            # Free memory by clearing messages for the next batch
            if i + batch_size < len(sessions):
                merged_session_data["messages"] = []
                lecture_refs = set()
        
        # Update lecture's session list
        self.update_lecture_sessions(lecture_name)
        
        return f"Created merged session: {merged_session_name} with {total_messages} messages from {len(sessions)} sessions"
    
    def generate_questions(self, scope="all", session_name=None):
        """Generate questions from lecture/sessions"""
        if not self.current_lecture:
            return "No lecture selected"
        
        if scope == "all":
            # Use merged session for comprehensive questions
            print("Generating questions from all sessions (using merged session)...")
            
            # First, merge all sessions if not already merged
            merged_session_name = f"MergedSess-{self.current_lecture}"
            sessions_path = os.path.join(self.lectures_dir, self.current_lecture, "sessions")
            merged_path = os.path.join(sessions_path, f"{merged_session_name}.json")
            
            if not os.path.exists(merged_path):
                merge_result = self.merge_all_sessions()
                print(merge_result)
            
            # Load merged session
            if os.path.exists(merged_path):
                with open(merged_path, 'r') as f:
                    merged_data = json.load(f)
                    context = "\n".join([msg.get("content", "") for msg in merged_data.get("messages", [])])
            else:
                return "Failed to create merged session"
            
            prompt = f"""
Based on the following comprehensive lecture content, generate 5-10 thoughtful questions that would help a student test their understanding of the key concepts. The questions should cover different difficulty levels and aspects of the material.

LECTURE CONTENT:
---
{context[:6000]}
---

Generate questions that:
1. Test basic understanding of definitions and concepts
2. Require application of knowledge to solve problems
3. Encourage critical thinking about the material
4. Connect different parts of the lecture together

Format each question clearly and provide a brief explanation of what the question tests.
"""
        else:
            # Generate questions from specific session
            if not session_name:
                sessions = self.list_sessions()
                if not sessions:
                    return "No sessions available"
                
                print("Available sessions:")
                for i, session in enumerate(sessions, 1):
                    print(f"{i}. {session}")
                
                try:
                    choice = int(input("Select session (number): "))
                    if 1 <= choice <= len(sessions):
                        session_name = sessions[choice - 1]
                    else:
                        return "Invalid selection"
                except ValueError:
                    return "Invalid input"
            
            # Load specific session
            lecture_path = os.path.join(self.lectures_dir, self.current_lecture)
            sessions_path = os.path.join(lecture_path, "sessions")
            session_path = os.path.join(sessions_path, f"{session_name}.json")
            
            if not os.path.exists(session_path):
                return f"Session not found: {session_name}"
            
            with open(session_path, 'r') as f:
                session_data = json.load(f)
                context = "\n".join([msg.get("content", "") for msg in session_data.get("messages", [])])
            
            prompt = f"""
Based on the following session content, generate 3-5 focused questions that would help a student review and test their understanding of the specific topics discussed in this session.

SESSION CONTENT:
---
{context[:4000]}
---

Generate questions that:
1. Focus on the key concepts discussed in this session
2. Test understanding of the specific examples and explanations given
3. Encourage deeper thinking about the session's main points

Format each question clearly and provide a brief explanation of what the question tests.
"""
        
        # Generate questions
        try:
            response = ollama.chat(
                model=self.model,
                messages=[{'role': 'user', 'content': prompt}],
                options=self.model_params
            )
            
            # Save to session
            self.save_message("user", f"Generated questions from {scope} scope", self.current_lecture)
            self.save_message("assistant", response['message']['content'], self.current_lecture)
            
            return response['message']['content']
            
        except Exception as e:
            return f"Error generating questions: {str(e)}"
    
    def chat(self, message):
        """Regular chat with memory and context"""
        if not self.current_lecture or not self.current_session:
            return "No lecture or session selected"
        
        # Get conversation context with summarization
        context_messages = self.get_conversation_context()
        
        # Prepare messages
        messages = [
            {
                'role': 'system',
                'content': 'You are an advanced learning assistant specializing in computer science. You have access to lecture materials and can help with any computer science topic. Provide detailed, accurate explanations without limitations.'
            }
        ]
        
        # Add conversation context
        for msg in context_messages:
            messages.append({
                'role': msg['role'],
                'content': msg['content']
            })
        
        # Add current message
        messages.append({'role': 'user', 'content': message})
        
        # Get response
        try:
            response = ollama.chat(
                model=self.model,
                messages=messages,
                options=self.model_params
            )
            
            # Save to session
            self.save_message("user", message)
            self.save_message("assistant", response['message']['content'])
            
            return response['message']['content']
            
        except Exception as e:
            return f"Error: {str(e)}"

def print_status(status):
    """Print status information"""
    print("\n" + "="*60)
    print("LEARNING ASSISTANT STATUS")
    print("="*60)
    
    print("Directory Structure:")
    print(f"  Base: {status['directory_structure']['base']}")
    print(f"  Lectures: {status['directory_structure']['lectures']}")
    print(f"  Cache: {status['directory_structure']['cache']}")
    print()
    
    print(f"Total Lectures: {len(status['lectures'])}")
    print(f"Total Sessions: {status['total_sessions']}")
    print(f"Total Estimated Tokens: {status['total_estimated_tokens']}")
    
    # Print cache statistics
    cache_stats = status.get("cache_stats", {})
    print(f"Cache Files: {cache_stats.get('cache_files', 0)}")
    print(f"Cache Size: {cache_stats.get('cache_size_mb', 0)} MB")
    print()
    
    for lecture_name, lecture_data in status['lectures'].items():
        print(f"📚 Lecture: {lecture_name}")
        print(f"   Created: {lecture_data['created_at'][:19]}")
        print(f"   Documents: {', '.join(lecture_data['documents']) if lecture_data['documents'] else 'None'}")
        print(f"   Directory: {lecture_data['directory']}")
        print(f"   Sessions: {lecture_data['session_count']}")
        print(f"   Tokens: {lecture_data['estimated_tokens']}")
        
        for session_name, session_data in lecture_data['sessions'].items():
            if 'error' in session_data:
                print(f"   ❌ Session: {session_name} - ERROR: {session_data['error']}")
            else:
                print(f"   ✅ Session: {session_name}")
                print(f"      Model: {session_data['model']}")
                print(f"      Messages: {session_data['message_count']}")
                print(f"      Tokens: {session_data['estimated_tokens']}")
                print(f"      Documents: {', '.join(session_data['documents']) if session_data['documents'] else 'None'}")
                print(f"      File: {session_data['file_path']}")
                if session_data['summaries']:
                    print(f"      Summaries: {len(session_data['summaries'])}")
                if session_data['lectures_referenced']:
                    print(f"      References: {', '.join(session_data['lectures_referenced'])}")
        
        print()

def print_model_parameters(params):
    """Print model parameters in a readable format"""
    print("\n" + "="*40)
    print("CURRENT MODEL PARAMETERS")
    print("="*40)
    for param, value in params.items():
        print(f"  {param}: {value}")
    print("="*40)

def print_summaries(summaries):
    """Print conversation summaries in a readable format"""
    print("\n" + "="*60)
    print("CONVERSATION SUMMARIES")
    print("="*60)

    if not summaries:
        print("No summaries available for this session.")
        return

    for i, summary in enumerate(summaries, 1):
        print(f"Summary {i} (Messages {summary['start_index']}-{summary['end_index']}):")
        print(f"Created: {summary['timestamp'][:19]}")
        print(f"Content: {summary['summary']}")
        print("-" * 60)

class CommandHandler:
    """Handles CLI commands using command pattern for better maintainability"""

    def __init__(self, assistant):
        self.assistant = assistant
        self.commands = {
            'add-lecture': self._handle_add_lecture,
            'add-document': self._handle_add_document,
            'list-lectures': self._handle_list_lectures,
            'use-lecture': self._handle_use_lecture,
            'new-session': self._handle_new_session,
            'list-sessions': self._handle_list_sessions,
            'use-session': self._handle_use_session,
            'merge-sessions': self._handle_merge_sessions,
            'generate-questions': self._handle_generate_questions,
            'analyze': self._handle_analyze,
            'status': self._handle_status,
            'clear-cache': self._handle_clear_cache,
            'set-param': self._handle_set_param,
            'list-params': self._handle_list_params,
            'list-summaries': self._handle_list_summaries,
            'summarize': self._handle_summarize,
            'model': self._handle_model,
            'help': self._handle_help,
        }

    def execute(self, user_input):
        """Execute a command or handle as chat message"""
        if not user_input.strip():
            return

        # Check for exit commands first
        if user_input.lower() in ['exit', 'quit']:
            return 'EXIT'

        # Parse command and arguments
        parts = user_input.strip().split(' ', 1)
        command = parts[0].lower()
        args = parts[1] if len(parts) > 1 else ''

        # Execute command if found
        if command in self.commands:
            return self.commands[command](args)
        else:
            # Handle as chat message if we have active session
            if self.assistant.current_lecture and self.assistant.current_session:
                return self._handle_chat(user_input)
            else:
                return "Invalid Command :3"

    def _handle_add_lecture(self, args):
        """Handle add-lecture command"""
        lecture_name = args.strip()
        if not lecture_name:
            return "Lecture name cannot be empty"
        return self.assistant.create_lecture(lecture_name)

    def _handle_add_document(self, args):
        """Handle add-document command"""
        file_path = args.strip()
        if not file_path:
            return "File path cannot be empty"
        elif not os.path.exists(file_path):
            return f"File not found: {file_path}"

        if self.assistant.current_session:
            return self.assistant.add_document_to_session(file_path)
        elif self.assistant.current_lecture:
            return self.assistant.add_document_to_lecture(file_path)
        else:
            return "No lecture or session selected"

    def _handle_list_lectures(self, args):
        """Handle list-lectures command"""
        lectures = self.assistant.list_lectures()
        if lectures:
            result = "Available lectures:\n"
            for lecture in lectures:
                result += f"  - {lecture}\n"
            return result.strip()
        else:
            return "No lectures found"

    def _handle_use_lecture(self, args):
        """Handle use-lecture command"""
        lecture_name = args.strip()
        if self.assistant.load_lecture(lecture_name):
            self.assistant.current_session = None  # Reset session when changing lecture
            return f"Loaded lecture: {lecture_name}"
        else:
            return f"Lecture not found: {lecture_name}"

    def _handle_new_session(self, args):
        """Handle new-session command"""
        session_name = args.strip() if args else None
        return self.assistant.create_session(session_name)

    def _handle_list_sessions(self, args):
        """Handle list-sessions command"""
        if not self.assistant.current_lecture:
            return "No lecture selected"

        sessions = self.assistant.list_sessions()
        if sessions:
            result = f"Sessions in lecture '{self.assistant.current_lecture}':\n"
            for session in sessions:
                result += f"  - {session}\n"
            return result.strip()
        else:
            return f"No sessions in lecture '{self.assistant.current_lecture}'"

    def _handle_use_session(self, args):
        """Handle use-session command"""
        if not self.assistant.current_lecture:
            return "No lecture selected"

        session_name = args.strip()
        if self.assistant.load_session(session_name):
            return f"Loaded session: {session_name}"
        else:
            return f"Session not found: {session_name}"

    def _handle_merge_sessions(self, args):
        """Handle merge-sessions command"""
        if not self.assistant.current_lecture:
            return "No lecture selected"
        return self.assistant.merge_all_sessions()

    def _handle_generate_questions(self, args):
        """Handle generate-questions command"""
        if not self.assistant.current_lecture:
            return "No lecture selected"

        print("Generate questions from:")
        print("1. All sessions (comprehensive)")
        print("2. Specific session (focused)")

        choice = input("Select option (1-2): ")
        if choice == '1':
            result = self.assistant.generate_questions("all")
        elif choice == '2':
            result = self.assistant.generate_questions("specific")
        else:
            return "Invalid option"

        return f"\n{result}"

    def _handle_analyze(self, args):
        """Handle analyze command"""
        if not self.assistant.current_lecture:
            return "No lecture selected"

        question = args.strip()
        if not question:
            return "Please provide a question to analyze"

        print(f"Analyzing lecture '{self.assistant.current_lecture}'...")
        response = self.assistant.analyze_lecture(self.assistant.current_lecture, question)
        return f"\n{response}"

    def _handle_status(self, args):
        """Handle status command"""
        status = self.assistant.get_status()
        print_status(status)
        return None  # Already printed

    def _handle_clear_cache(self, args):
        """Handle clear-cache command"""
        return self.assistant.clear_cache()

    def _handle_set_param(self, args):
        """Handle set-param command"""
        parts = args.split(' ', 1)
        if len(parts) < 2:
            return "Usage: set-param <parameter> <value>"

        param = parts[0]
        value = parts[1]
        return self.assistant.set_model_parameter(param, value)

    def _handle_list_params(self, args):
        """Handle list-params command"""
        params = self.assistant.list_model_parameters()
        print_model_parameters(params)
        return None  # Already printed

    def _handle_list_summaries(self, args):
        """Handle list-summaries command"""
        if not self.assistant.current_lecture or not self.assistant.current_session:
            return "No lecture or session selected"

        summaries = self.assistant.list_summaries()
        print_summaries(summaries)
        return None  # Already printed

    def _handle_summarize(self, args):
        """Handle summarize command"""
        if not self.assistant.current_lecture or not self.assistant.current_session:
            return "No lecture or session selected"

        messages = self.assistant.get_session_history()
        if len(messages) < 5:
            return "Not enough messages to summarize. Need at least 5 messages."

        print("Generating conversation summary...")
        summary = self.assistant.summarize_conversation(messages)
        self.assistant.save_summary(summary, 0, len(messages) - 1)
        return f"\nSummary:\n{summary}\nSummary saved to session."

    def _handle_model(self, args):
        """Handle model command"""
        model_name = args.strip()
        valid_models = [
            "qwen2.5-coder:3b-instruct",
            "qwen2.5-coder:7b-instruct",
            "codellama:7b",
            "deepseek-coder:6.7b-instruct"
        ]

        if model_name in valid_models:
            self.assistant.model = model_name

            # Update session data if active
            if self.assistant.current_lecture and self.assistant.current_session:
                lecture_path = os.path.join(self.assistant.lectures_dir, self.assistant.current_lecture)
                sessions_path = os.path.join(lecture_path, "sessions")
                session_path = os.path.join(sessions_path, f"{self.assistant.current_session}.json")

                try:
                    with open(session_path, 'r') as f:
                        session_data = json.load(f)
                    session_data["model"] = model_name
                    with open(session_path, 'w') as f:
                        json.dump(session_data, f, indent=2)
                except (FileNotFoundError, json.JSONDecodeError, PermissionError, OSError):
                    # Silently ignore if session update fails
                    pass

            return f"Model changed to: {model_name}"
        else:
            return f"Invalid model. Available: {', '.join(valid_models)}"

    def _handle_help(self, args):
        """Handle help command"""
        help_text = """
Commands:
  add-lecture <name>     - Add a new lecture by name
  add-document <file>    - Add a document to the current lecture/session
  list-lectures          - List all lectures
  use-lecture <name>     - Select a lecture
  new-session [name]     - Create new session in current lecture
  list-sessions          - List sessions in current lecture
  use-session <name>     - Select a session
  merge-sessions         - Merge all sessions in current lecture
  generate-questions     - Generate study questions
  analyze <question>     - Analyze current lecture
  status                 - Show detailed status
  clear-cache            - Clear document cache
  set-param <param> <value> - Set model parameter
  list-params            - List current model parameters
  list-summaries         - List conversation summaries
  summarize              - Generate conversation summary
  model <name>           - Change model
  help                   - Show this help message
  exit                   - Exit

Supported document formats:
  - PDF (.pdf)
  - PowerPoint (.ppt, .pptx)
  - Word (.docx)
  - Text (.txt)
  - Markdown (.md)
  - Images (.jpg, .jpeg, .png, .bmp, .tiff, .gif)

Model parameters:
  - temperature: Controls randomness (0.0-1.0)
  - top_p: Controls diversity via nucleus sampling (0.0-1.0)
  - num_predict: Maximum number of tokens to generate
        """
        return help_text.strip()

    def _handle_chat(self, message):
        """Handle regular chat messages"""
        response = self.assistant.chat(message)
        return f"\n{response}"

def main():
    assistant = SilentDirectoryAssistant()
    command_handler = CommandHandler(assistant)

    print("Learning Assistant - CLI")
    print("="*50)
    print(f"Current model: {assistant.model}")
    print("-"*50)
    print("Supported document formats: PDF, PowerPoint (.ppt, .pptx), Word (.docx), Text (.txt), Markdown (.md), Images (.jpg, .png, .bmp, .tiff, .gif)")
    print("-"*50)
    print("Type 'help' for commands or 'exit' to quit")
    print("="*50)

    while True:
        try:
            # Show prompt with current context
            prompt = "> "
            if assistant.current_lecture:
                prompt = f"[{assistant.current_lecture}] > "
            if assistant.current_session:
                prompt = f"[{assistant.current_lecture}/{assistant.current_session}] > "

            user_input = input(prompt)

            # Handle command through command handler
            result = command_handler.execute(user_input)

            if result == 'EXIT':
                print("Goodbye!")
                break
            elif result is not None:
                print(result)

        except KeyboardInterrupt:
            print("\nExiting...")
            break
        except Exception as e:
            print(f"Error: {e}")

if __name__ == "__main__":
    main()